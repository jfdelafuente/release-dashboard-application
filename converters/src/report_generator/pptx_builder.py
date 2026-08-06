"""Ensamblado del informe .pptx: portada, tarjetas de KPI y diapositivas de gráfica.

Estilo inspirado en el dashboard de postmortem y en el informe manual de
referencia ("PostMortem-PostProducción ... .pptx"), pero construido desde
cero con python-pptx (no se parte de una plantilla corporativa — ver
research.md §1).
"""
import io
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from report_generator.chart_utils import (
    COLOR_ORANGE, COLOR_INK, COLOR_INK_LIGHT, COLOR_BORDER, COLOR_SUCCESS, COLOR_DANGER,
)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

_RGB_ORANGE = RGBColor.from_string(COLOR_ORANGE.lstrip("#"))
_RGB_INK = RGBColor.from_string(COLOR_INK.lstrip("#"))
_RGB_INK_LIGHT = RGBColor.from_string(COLOR_INK_LIGHT.lstrip("#"))
_RGB_BORDER = RGBColor.from_string(COLOR_BORDER.lstrip("#"))
_RGB_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_RGB_SUCCESS = RGBColor.from_string(COLOR_SUCCESS.lstrip("#"))
_RGB_DANGER = RGBColor.from_string(COLOR_DANGER.lstrip("#"))

# Paleta del "hero" del rediseño MASORANGE (fondo negro), distinta de la
# usada en las tarjetas de KPI (fondo blanco) — ver
# "MASORANGE dashboard redesign/Release Dashboard.dc.html", sección Portal.
_RGB_HERO_SUBTITLE = RGBColor.from_string("B8B2A9")
_RGB_HERO_DEPT_LABEL = RGBColor.from_string("9E988E")

# Ruta relativa al propio módulo (no al cwd): sigue siendo válida aunque
# generate_postmortem_report.py haga chdir a otro repo (ver _maybe_chdir).
_LOGO_PATH = Path(__file__).parent / "assets" / "masorange-logo-positive.png"


def _blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _add_textbox(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def _add_title_bar(slide, title_text):
    _add_textbox(
        slide, Inches(0.5), Inches(0.25), SLIDE_WIDTH - Inches(1), Inches(0.6),
        title_text, size=24, color=_RGB_INK, bold=True,
    )
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.85), SLIDE_WIDTH - Inches(1), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = _RGB_ORANGE
    line.line.fill.background()


def new_presentation(release_name):
    """Crea la Presentation con la portada.

    Sigue el "hero" de portal del rediseño MASORANGE (fondo negro, logo,
    eyebrow naranja, titular grande en blanco, subtítulo gris) — ver
    "MASORANGE dashboard redesign/Release Dashboard.dc.html".
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide = _blank_slide(prs)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    background.fill.solid()
    background.fill.fore_color.rgb = _RGB_INK
    background.line.fill.background()
    background.shadow.inherit = False

    if _LOGO_PATH.exists():
        slide.shapes.add_picture(str(_LOGO_PATH), Inches(0.6), Inches(0.55), height=Inches(0.32))

    _add_textbox(
        slide, Inches(2.1), Inches(0.5), Inches(4), Inches(0.4),
        "Customer & Service Operations", size=11, color=_RGB_HERO_DEPT_LABEL, bold=True,
    )

    _add_textbox(
        slide, Inches(0.6), Inches(2.9), SLIDE_WIDTH - Inches(1.2), Inches(0.4),
        "POSTMORTEM · MASORANGE", size=13, color=_RGB_ORANGE, bold=True,
    )
    _add_textbox(
        slide, Inches(0.55), Inches(3.3), SLIDE_WIDTH - Inches(1.2), Inches(1.3),
        release_name, size=54, color=_RGB_WHITE, bold=True,
    )
    _add_textbox(
        slide, Inches(0.6), Inches(4.65), SLIDE_WIDTH - Inches(2.5), Inches(0.7),
        "Informe de postmortem — rendimiento de incidencias por despliegue y evolución operativa.",
        size=14, color=_RGB_HERO_SUBTITLE,
    )
    return prs


def _add_kpi_card(slide, left, top, width, height, value_text, label_text, detail_lines=None, value_color=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = _RGB_WHITE
    card.line.color.rgb = _RGB_BORDER
    card.line.width = Pt(1)
    card.shadow.inherit = False

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(4), height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = _RGB_ORANGE
    stripe.line.fill.background()

    pad = Inches(0.2)
    _add_textbox(slide, left + pad, top + Inches(0.15), width - 2 * pad, Inches(0.55), value_text, size=28, color=value_color or _RGB_INK, bold=True)
    _add_textbox(slide, left + pad, top + Inches(0.7), width - 2 * pad, Inches(0.3), label_text, size=11, color=_RGB_INK_LIGHT, bold=True)
    if detail_lines:
        box = slide.shapes.add_textbox(left + pad, top + Inches(1.0), width - 2 * pad, height - Inches(1.1))
        tf = box.text_frame
        tf.word_wrap = True
        for idx, line in enumerate(detail_lines):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(9)
            run.font.color.rgb = _RGB_INK_LIGHT
            run.font.name = "Calibri"


def _kpi_status_color(pct, target_pct):
    """Verde si `pct` alcanza el objetivo, rojo si queda por debajo."""
    return _RGB_SUCCESS if pct >= target_pct else _RGB_DANGER


def add_kpi_and_chart_slide(prs, release, chart_png, target_pct):
    """Diapositiva "Métricas Globales": 3 tarjetas de KPI (Total Incidencias,
    % Resueltas PaP, % Resueltas Mesa — estas dos en verde/rojo según
    `target_pct`) junto a la gráfica "Incidencias por Release".

    `release` es un dict de release_kpis_data.build_releases()/find_release()
    — todos los valores vienen de releases-data.js, no del JSON de postmortem.
    """
    slide = _blank_slide(prs)
    _add_title_bar(slide, f"Métricas Globales — {release['name']}")

    card_left = Inches(0.5)
    card_width = Inches(3.6)
    card_height = Inches(1.7)
    card_gap = Inches(0.25)
    top0 = Inches(1.3)

    _add_kpi_card(
        slide, card_left, top0, card_width, card_height,
        str(release["total_incidencias"]), "Total Incidencias",
    )

    pap_top = top0 + card_height + card_gap
    _add_kpi_card(
        slide, card_left, pap_top, card_width, card_height,
        f"{release['pct_pap']}%", "% Resueltas PaP",
        detail_lines=[
            f"{release['pap_resueltas']} de {release['pap_entrada']} incidencias PaP resueltas el día del PaP",
            f"Objetivo: {target_pct}%",
        ],
        value_color=_kpi_status_color(release["pct_pap"], target_pct),
    )

    mesa_top = pap_top + card_height + card_gap
    _add_kpi_card(
        slide, card_left, mesa_top, card_width, card_height,
        f"{release['pct_first_week']}%", "% Resueltas Mesa",
        detail_lines=[
            f"{release['post_resueltas']} de {release['post_entrada']} incidencias Mesa",
            f"Objetivo: {target_pct}%",
        ],
        value_color=_kpi_status_color(release["pct_first_week"], target_pct),
    )

    chart_left = card_left + card_width + Inches(0.4)
    chart_width = SLIDE_WIDTH - chart_left - Inches(0.5)
    slide.shapes.add_picture(io.BytesIO(chart_png), chart_left, top0, width=chart_width)
    return slide


def add_dual_chart_slide(prs, title, charts):
    """Diapositiva final con las gráficas de `charts` ([(título, png_bytes), ...])
    repartidas en columnas iguales, cada una con su propio título."""
    slide = _blank_slide(prs)
    _add_title_bar(slide, title)

    margin = Inches(0.5)
    gap = Inches(0.4)
    n = len(charts)
    col_width = (SLIDE_WIDTH - 2 * margin - (n - 1) * gap) / n
    top = Inches(1.3)
    label_height = Inches(0.4)

    for idx, (chart_title, png_bytes) in enumerate(charts):
        left = margin + idx * (col_width + gap)
        _add_textbox(slide, left, top, col_width, label_height, chart_title, size=16, color=_RGB_INK, bold=True, align=PP_ALIGN.CENTER)
        slide.shapes.add_picture(io.BytesIO(png_bytes), left, top + label_height + Inches(0.1), width=col_width)

    return slide
