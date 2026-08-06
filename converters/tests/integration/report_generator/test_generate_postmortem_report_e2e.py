"""Test end-to-end: dataset sintético -> generate_report() -> abrir el .pptx resultante."""
import json

from pptx import Presentation

from generate_postmortem_report import generate_report
from report_generator.kpi_calculator import calculate_kpis


def _write_release(output_dir, release_name, records):
    path = output_dir / f"{release_name}-postmortem.json"
    path.write_text(
        json.dumps({"_metadata": {"release_name": release_name}, "data": records}, ensure_ascii=False),
        encoding="utf-8",
    )
    return path


def _synthetic_records():
    return [
        {"ID de incidencia": "INC1", "Estatus": "Cerrado", "Despliegue": "PAP",
         "Fecha de envío": "10/06/2026 08:00", "Fecha de última resolución": "10/06/2026 10:00",
         "Grupo asignado": "SOP_A"},
        {"ID de incidencia": "INC2", "Estatus": "Asignado", "Despliegue": "PAP",
         "Fecha de envío": "10/06/2026 09:00", "Fecha de última resolución": "",
         "Grupo asignado": "SOP_B"},
        {"ID de incidencia": "INC3", "Estatus": "Cerrado", "Despliegue": "MESA",
         "Fecha de envío": "12/06/2026 08:00", "Fecha de última resolución": "12/06/2026 09:30",
         "Grupo asignado": "SOP_A"},
        {"ID de incidencia": "INC4", "Estatus": "Pendiente", "Despliegue": "MESA",
         "Fecha de envío": "12/06/2026 09:00", "Fecha de última resolución": "",
         "Grupo asignado": "SOP_B"},
    ]


class TestGenerateReportE2E:
    def test_generates_pptx_with_expected_slides_and_kpis(self, tmp_path, monkeypatch):
        output_dir = tmp_path / "data" / "output"
        output_dir.mkdir(parents=True)
        records = _synthetic_records()
        source_path = _write_release(output_dir, "2026TEST", records)

        monkeypatch.setattr(
            "generate_postmortem_report.load_postmortem_records",
            lambda release_name, output_dir=None: records if release_name == "2026TEST" else None,
        )
        monkeypatch.setattr(
            "generate_postmortem_report.find_postmortem_file",
            lambda release_name, output_dir=None: source_path if release_name == "2026TEST" else None,
        )

        report_path = tmp_path / "report.pptx"
        result = generate_report("2026TEST", output_path=report_path)

        assert result["success"] is True
        assert report_path.exists()

        prs = Presentation(str(report_path))
        # Portada + KPIs (las gráficas propias de postmortem ya no se
        # incluyen en el informe; releases-data.js no está disponible desde
        # el cwd de test, así que tampoco hay slides de contexto de release-kpis)
        assert len(prs.slides) == 2

        expected_kpis = calculate_kpis(records)
        kpi_slide = prs.slides[1]
        slide_text = "\n".join(
            s.text_frame.text for s in kpi_slide.shapes if s.has_text_frame
        )
        assert str(expected_kpis["total_incidencias"]) in slide_text
        assert f"{expected_kpis['pct_cerradas']}%" in slide_text

    def test_unknown_release_returns_error_without_creating_file(self, tmp_path, monkeypatch):
        from report_generator.data_loader import ReleaseNotFoundError

        def _raise(*args, **kwargs):
            raise ReleaseNotFoundError("No hay datos de postmortem cargados para la release 'NOPE'")

        monkeypatch.setattr("generate_postmortem_report.load_postmortem_records", _raise)

        report_path = tmp_path / "report.pptx"
        result = generate_report("NOPE", output_path=report_path)

        assert result["success"] is False
        assert "NOPE" in result["error"]
        assert not report_path.exists()

    def test_locked_output_file_raises_friendly_message(self, tmp_path, monkeypatch):
        """Si el .pptx anterior está abierto en otro programa (p. ej.
        PowerPoint), Windows deniega la escritura con PermissionError. El
        mensaje debe ser accionable, no el error crudo del sistema."""
        records = _synthetic_records()
        monkeypatch.setattr(
            "generate_postmortem_report.load_postmortem_records",
            lambda release_name, output_dir=None: records,
        )
        monkeypatch.setattr(
            "generate_postmortem_report.find_postmortem_file",
            lambda release_name, output_dir=None: tmp_path / "2026TEST-postmortem.json",
        )
        monkeypatch.setattr(
            "pptx.presentation.Presentation.save",
            lambda self, path: (_ for _ in ()).throw(PermissionError("[Errno 13] Permission denied")),
        )

        report_path = tmp_path / "2026TEST-postmortem-report.pptx"
        try:
            generate_report("2026TEST", output_path=report_path)
            assert False, "Se esperaba PermissionError"
        except PermissionError as e:
            assert "abierto en otro programa" in str(e)
            assert report_path.name in str(e)


class TestGenerateReportCaching:
    """El informe no debe regenerarse si ya existe y es más reciente que sus
    fuentes de datos — evita repetir el renderizado de las 7 gráficas en
    cada clic cuando nadie ha subido datos nuevos desde la última vez."""

    def test_skips_regeneration_when_report_is_newer_than_source(self, tmp_path, monkeypatch):
        import time

        output_dir = tmp_path / "data" / "output"
        output_dir.mkdir(parents=True)
        records = _synthetic_records()
        source_path = _write_release(output_dir, "2026TEST", records)

        monkeypatch.setattr(
            "generate_postmortem_report.find_postmortem_file",
            lambda release_name, output_dir=None: source_path,
        )
        calls = {"count": 0}

        def _counting_load(release_name, output_dir=None):
            calls["count"] += 1
            return records

        monkeypatch.setattr("generate_postmortem_report.load_postmortem_records", _counting_load)

        report_path = tmp_path / "report.pptx"
        first = generate_report("2026TEST", output_path=report_path)
        assert first["success"] is True
        assert calls["count"] == 1
        first_mtime = report_path.stat().st_mtime

        # Segunda llamada sin cambios en la fuente: no debe regenerar.
        second = generate_report("2026TEST", output_path=report_path)
        assert second["success"] is True
        assert calls["count"] == 1  # sigue en 1: no se volvió a llamar a load_postmortem_records
        assert report_path.stat().st_mtime == first_mtime

        # La fuente cambia (más reciente que el informe): ahora sí debe regenerar.
        time.sleep(0.05)
        source_path.write_text(source_path.read_text(encoding="utf-8"), encoding="utf-8")
        assert source_path.stat().st_mtime > first_mtime

        third = generate_report("2026TEST", output_path=report_path)
        assert third["success"] is True
        assert calls["count"] == 2
