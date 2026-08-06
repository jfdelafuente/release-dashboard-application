"""Test de integración: el informe de una release incluye el contexto
comparativo de release-kpis con TODAS las releases, no solo la seleccionada."""
import json

from pptx import Presentation

from generate_postmortem_report import generate_report


_RELEASES_JS = """"use strict";
const RAW_RELEASES = [
  ["2026R6", 2026, "7-jun.", "Junio", 53, 46, 38, 33],
  ["2026R7", 2026, "7-jul.", "Julio", 70, 54, 107, 84],
  ["2026R8", 2026, "7-ago.", "Agosto", 40, 35, 20, 18],
];
"""


def _write_release(output_dir, release_name, records):
    path = output_dir / f"{release_name}-postmortem.json"
    path.write_text(
        json.dumps({"_metadata": {"release_name": release_name}, "data": records}, ensure_ascii=False),
        encoding="utf-8",
    )


def _minimal_records():
    return [
        {"ID de incidencia": "INC1", "Estatus": "Cerrado", "Despliegue": "MESA",
         "Fecha de envío": "10/06/2026 08:00", "Fecha de última resolución": "10/06/2026 09:00",
         "Grupo asignado": "SOP_A"},
        {"ID de incidencia": "INC2", "Estatus": "Cerrado", "Despliegue": "PAP",
         "Fecha de envío": "10/06/2026 09:00", "Fecha de última resolución": "10/06/2026 10:00",
         "Grupo asignado": "SOP_B"},
    ]


class TestReleaseKpisContextInReport:
    def test_report_includes_all_three_releases_not_just_selected(self, tmp_path, monkeypatch):
        releases_data_path = tmp_path / "releases-data.js"
        releases_data_path.write_text(_RELEASES_JS, encoding="utf-8")

        records = _minimal_records()
        monkeypatch.setattr(
            "generate_postmortem_report.load_postmortem_records",
            lambda release_name, output_dir=None: records if release_name == "2026R6" else None,
        )
        monkeypatch.setattr(
            "generate_postmortem_report.find_postmortem_file",
            lambda release_name, output_dir=None: tmp_path / f"{release_name}-postmortem.json",
        )
        monkeypatch.setattr(
            "generate_postmortem_report.load_release_kpis_context",
            lambda path=None: __import__("report_generator.release_kpis_data", fromlist=["load_release_kpis_context"]).load_release_kpis_context(path=releases_data_path),
        )

        report_path = tmp_path / "report.pptx"
        result = generate_report("2026R6", output_path=report_path)
        assert result["success"] is True

        prs = Presentation(str(report_path))
        # Portada + KPIs + 3 release-kpis = 5
        assert len(prs.slides) == 5

        titles = [
            s.text_frame.text
            for slide in prs.slides
            for s in slide.shapes
            if s.has_text_frame and s.text_frame.text.strip()
        ]
        assert any("Incidencias por Release" in t for t in titles)
        assert any("KPI % PaP" in t for t in titles)
        assert any("KPI % 1ª semana" in t for t in titles)

    def test_report_still_generates_when_release_kpis_data_missing(self, tmp_path, monkeypatch):
        """Ver Edge Cases de spec.md: si falta la sección de release-kpis, el
        informe se genera igualmente con los KPIs de la release."""
        records = _minimal_records()
        monkeypatch.setattr(
            "generate_postmortem_report.load_postmortem_records",
            lambda release_name, output_dir=None: records,
        )
        monkeypatch.setattr(
            "generate_postmortem_report.find_postmortem_file",
            lambda release_name, output_dir=None: tmp_path / f"{release_name}-postmortem.json",
        )

        def _raise(path=None):
            raise FileNotFoundError("releases-data.js no encontrado")

        monkeypatch.setattr("generate_postmortem_report.load_release_kpis_context", _raise)

        report_path = tmp_path / "report.pptx"
        result = generate_report("2026R6", output_path=report_path)
        assert result["success"] is True

        prs = Presentation(str(report_path))
        # Portada + KPIs, sin contexto de release-kpis
        assert len(prs.slides) == 2
